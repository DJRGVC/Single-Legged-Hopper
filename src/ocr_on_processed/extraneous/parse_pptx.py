from pptx import Presentation
import os
import shutil
import pandas as pd

# Define directories to store images and text
base_dir = "../raw_data"
images_dir = os.path.join(base_dir, "images")
labels_file = os.path.join(base_dir, "labels.csv")  # CSV for PyTorch-friendly format

# Remove existing directories and files if needed
if os.path.exists(base_dir):
    shutil.rmtree(base_dir)
os.makedirs(images_dir)

# Load the presentation
prs = Presentation("../presentation_data/RIC.pptx")  # Change this to the path of your downloaded presentation

# List to store slide data (image paths and corresponding text)
slide_data = []

# Iterate over slides
for idx, slide in enumerate(prs.slides):
    slide_text = ""
    image_idx = 1
    
    # Extract text
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            slide_text += shape.text.strip() + " "
    
    # Clean up slide text (for CSV)
    slide_text = slide_text.strip().replace("\n", " ")
    
    # Extract images
    for shape in slide.shapes:
        if hasattr(shape, "image"):
            image = shape.image
            image_filename = f"slide_{idx + 1}_image_{image_idx}.png"
            image_file_path = os.path.join(images_dir, image_filename)
            
            # Save the image
            with open(image_file_path, "wb") as img_out:
                img_out.write(image.blob)
            
            # Append image path and corresponding text to the list
            slide_data.append({
                "image_path": image_file_path,
                "text": slide_text
            })
            
            image_idx += 1

# Convert the list to a pandas DataFrame for easy saving as CSV
df = pd.DataFrame(slide_data, columns=["image_path", "text"])

# Save the DataFrame to a CSV file (this will act as the labels file for PyTorch)
df.to_csv(labels_file, index=False)

print(f"Images saved to {images_dir}")
print(f"Labels saved to {labels_file}")

